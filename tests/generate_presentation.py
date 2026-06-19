from pptx import Presentation
from pptx.util import Inches

def create_presentation():
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Fact-Check Agent Web App"
    subtitle.text = "Automated Document Verification System\nArchitecture and Strategy"
    
    # Architecture Slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "System Architecture"
    
    tf = body_shape.text_frame
    tf.text = "Frontend: Streamlit Web UI"
    p = tf.add_paragraph()
    p.text = "Backend: FastAPI Service"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "LLM Engine: Gemini 1.5 Pro for extraction & verification logic"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Web Search: Tavily API for precise fact retrieval"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Database: PostgreSQL (via SQLAlchemy) for job and claim tracking"
    p.level = 0
    
    # Evaluation Logic Slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.shapes.placeholders[1]
    title_shape.text = "Evaluation Engine & Trap Logic"
    
    tf = body_shape.text_frame
    tf.text = "Handling Trap Documents:"
    p = tf.add_paragraph()
    p.text = "Differentiates between FALSE (hallucinated) and INACCURATE (outdated)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Extracts explicit date context from claims to check against publication dates of evidence."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Uses credibility scoring to weigh official sources over blogs."
    p.level = 1
    
    prs.save("Fact_Check_Agent_Architecture.pptx")
    print("Created Fact_Check_Agent_Architecture.pptx")

if __name__ == "__main__":
    create_presentation()
